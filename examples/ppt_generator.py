#!/usr/bin/env python3

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Any, Optional


class PPTGenerator:
    def __init__(self):
        self.presentation = Presentation()
        self.shapes_data = []
        self.layouts_data = []
        self.theme_data = {}
        
    def load_json_files(self, shapes_file: str, layouts_file: str, theme_file: str):
        """Load data from JSON files"""
        try:
            with open(shapes_file, 'r', encoding='utf-8') as f:
                self.shapes_data = json.load(f)
            
            with open(layouts_file, 'r', encoding='utf-8') as f:
                self.layouts_data = json.load(f)
                
            with open(theme_file, 'r', encoding='utf-8') as f:
                self.theme_data = json.load(f)
                
            print(f"Loaded {len(self.shapes_data)} slide(s) with shapes")
            print(f"Loaded {len(self.layouts_data)} layout(s)")
            print(f"Loaded theme: {self.theme_data.get('theme_name', 'Unknown')}")
            
        except Exception as e:
            raise Exception(f"Error loading JSON files: {str(e)}")
    
    def emu_to_inches(self, emu_value: int) -> float:
        """Convert EMU (English Metric Units) to inches"""
        return emu_value / 914400.0
    
    def parse_shape_type_number(self, shape_type_str: str) -> Optional[int]:
        """Extract shape type number from string like 'FREEFORM (5)' or 'FREEFORM (5) (5)'"""
        import re
        # Handle both single and duplicated number formats
        matches = re.findall(r'\((\d+)\)', shape_type_str)
        if matches:
            # Take the first number found
            return int(matches[0])
        return None
    
    def get_mso_auto_shape_from_name(self, shape_name: str):
        """Convert auto shape name to MSO_SHAPE constant"""
        from pptx.enum.shapes import MSO_SHAPE
        
        # Clean up the shape name
        shape_name = shape_name.upper().strip()
        
        # Comprehensive mapping of MSO_AUTO_SHAPE_TYPE names to MSO_SHAPE constants
        auto_shape_mapping = {
            # Action Buttons
            'ACTION_BUTTON_BACK_OR_PREVIOUS': MSO_SHAPE.ACTION_BUTTON_BACK_OR_PREVIOUS,
            'ACTION_BUTTON_BEGINNING': MSO_SHAPE.ACTION_BUTTON_BEGINNING,
            'ACTION_BUTTON_CUSTOM': MSO_SHAPE.ACTION_BUTTON_CUSTOM,
            'ACTION_BUTTON_DOCUMENT': MSO_SHAPE.ACTION_BUTTON_DOCUMENT,
            'ACTION_BUTTON_END': MSO_SHAPE.ACTION_BUTTON_END,
            'ACTION_BUTTON_FORWARD_OR_NEXT': MSO_SHAPE.ACTION_BUTTON_FORWARD_OR_NEXT,
            'ACTION_BUTTON_HELP': MSO_SHAPE.ACTION_BUTTON_HELP,
            'ACTION_BUTTON_HOME': MSO_SHAPE.ACTION_BUTTON_HOME,
            'ACTION_BUTTON_INFORMATION': MSO_SHAPE.ACTION_BUTTON_INFORMATION,
            'ACTION_BUTTON_MOVIE': MSO_SHAPE.ACTION_BUTTON_MOVIE,
            'ACTION_BUTTON_RETURN': MSO_SHAPE.ACTION_BUTTON_RETURN,
            'ACTION_BUTTON_SOUND': MSO_SHAPE.ACTION_BUTTON_SOUND,
            
            # Basic Shapes
            'ARC': MSO_SHAPE.ARC,
            'BALLOON': MSO_SHAPE.BALLOON,
            'BENT_ARROW': MSO_SHAPE.BENT_ARROW,
            'BENT_UP_ARROW': MSO_SHAPE.BENT_UP_ARROW,
            'BEVEL': MSO_SHAPE.BEVEL,
            'BLOCK_ARC': MSO_SHAPE.BLOCK_ARC,
            'CAN': MSO_SHAPE.CAN,
            'CHEVRON': MSO_SHAPE.CHEVRON,
            'CHORD': MSO_SHAPE.CHORD,
            'CIRCULAR_ARROW': MSO_SHAPE.CIRCULAR_ARROW,
            'CLOUD': MSO_SHAPE.CLOUD,
            'CLOUD_CALLOUT': MSO_SHAPE.CLOUD_CALLOUT,
            'CORNER': MSO_SHAPE.CORNER,
            'CORNER_TABS': MSO_SHAPE.CORNER_TABS,
            'CROSS': MSO_SHAPE.CROSS,
            'CUBE': MSO_SHAPE.CUBE,
            
            # Curved Arrows
            'CURVED_DOWN_ARROW': MSO_SHAPE.CURVED_DOWN_ARROW,
            'CURVED_DOWN_RIBBON': MSO_SHAPE.CURVED_DOWN_RIBBON,
            'CURVED_LEFT_ARROW': MSO_SHAPE.CURVED_LEFT_ARROW,
            'CURVED_RIGHT_ARROW': MSO_SHAPE.CURVED_RIGHT_ARROW,
            'CURVED_UP_ARROW': MSO_SHAPE.CURVED_UP_ARROW,
            'CURVED_UP_RIBBON': MSO_SHAPE.CURVED_UP_RIBBON,
            
            # Polygons
            'DECAGON': MSO_SHAPE.DECAGON,
            'DIAMOND': MSO_SHAPE.DIAMOND,
            'DODECAGON': MSO_SHAPE.DODECAGON,
            'DONUT': MSO_SHAPE.DONUT,
            'DOUBLE_BRACE': MSO_SHAPE.DOUBLE_BRACE,
            'DOUBLE_BRACKET': MSO_SHAPE.DOUBLE_BRACKET,
            'DOUBLE_WAVE': MSO_SHAPE.DOUBLE_WAVE,
            
            # Arrows
            'DOWN_ARROW': MSO_SHAPE.DOWN_ARROW,
            'DOWN_ARROW_CALLOUT': MSO_SHAPE.DOWN_ARROW_CALLOUT,
            'DOWN_RIBBON': MSO_SHAPE.DOWN_RIBBON,
            
            # Explosions
            'EXPLOSION1': MSO_SHAPE.EXPLOSION1,
            'EXPLOSION2': MSO_SHAPE.EXPLOSION2,
            
            # Flowchart Shapes
            'FLOWCHART_ALTERNATE_PROCESS': MSO_SHAPE.FLOWCHART_ALTERNATE_PROCESS,
            'FLOWCHART_CARD': MSO_SHAPE.FLOWCHART_CARD,
            'FLOWCHART_COLLATE': MSO_SHAPE.FLOWCHART_COLLATE,
            'FLOWCHART_CONNECTOR': MSO_SHAPE.FLOWCHART_CONNECTOR,
            'FLOWCHART_DATA': MSO_SHAPE.FLOWCHART_DATA,
            'FLOWCHART_DECISION': MSO_SHAPE.FLOWCHART_DECISION,
            'FLOWCHART_DELAY': MSO_SHAPE.FLOWCHART_DELAY,
            'FLOWCHART_DIRECT_ACCESS_STORAGE': MSO_SHAPE.FLOWCHART_DIRECT_ACCESS_STORAGE,
            'FLOWCHART_DISPLAY': MSO_SHAPE.FLOWCHART_DISPLAY,
            'FLOWCHART_DOCUMENT': MSO_SHAPE.FLOWCHART_DOCUMENT,
            'FLOWCHART_EXTRACT': MSO_SHAPE.FLOWCHART_EXTRACT,
            'FLOWCHART_INTERNAL_STORAGE': MSO_SHAPE.FLOWCHART_INTERNAL_STORAGE,
            'FLOWCHART_MAGNETIC_DISK': MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
            'FLOWCHART_MANUAL_INPUT': MSO_SHAPE.FLOWCHART_MANUAL_INPUT,
            'FLOWCHART_MANUAL_OPERATION': MSO_SHAPE.FLOWCHART_MANUAL_OPERATION,
            'FLOWCHART_MERGE': MSO_SHAPE.FLOWCHART_MERGE,
            'FLOWCHART_MULTIDOCUMENT': MSO_SHAPE.FLOWCHART_MULTIDOCUMENT,
            'FLOWCHART_OFFLINE_STORAGE': MSO_SHAPE.FLOWCHART_OFFLINE_STORAGE,
            'FLOWCHART_OFFPAGE_CONNECTOR': MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR,
            'FLOWCHART_OR': MSO_SHAPE.FLOWCHART_OR,
            'FLOWCHART_PREDEFINED_PROCESS': MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
            'FLOWCHART_PREPARATION': MSO_SHAPE.FLOWCHART_PREPARATION,
            'FLOWCHART_PROCESS': MSO_SHAPE.FLOWCHART_PROCESS,
            'FLOWCHART_PUNCHED_TAPE': MSO_SHAPE.FLOWCHART_PUNCHED_TAPE,
            'FLOWCHART_SEQUENTIAL_ACCESS_STORAGE': MSO_SHAPE.FLOWCHART_SEQUENTIAL_ACCESS_STORAGE,
            'FLOWCHART_SORT': MSO_SHAPE.FLOWCHART_SORT,
            'FLOWCHART_STORED_DATA': MSO_SHAPE.FLOWCHART_STORED_DATA,
            'FLOWCHART_SUMMING_JUNCTION': MSO_SHAPE.FLOWCHART_SUMMING_JUNCTION,
            'FLOWCHART_TERMINATOR': MSO_SHAPE.FLOWCHART_TERMINATOR,
            
            # Other Shapes
            'FOLDED_CORNER': MSO_SHAPE.FOLDED_CORNER,
            'FRAME': MSO_SHAPE.FRAME,
            'FUNNEL': MSO_SHAPE.FUNNEL,
            'GEAR_6': MSO_SHAPE.GEAR_6,
            'GEAR_9': MSO_SHAPE.GEAR_9,
            'HALF_FRAME': MSO_SHAPE.HALF_FRAME,
            'HEART': MSO_SHAPE.HEART,
            'HEPTAGON': MSO_SHAPE.HEPTAGON,
            'HEXAGON': MSO_SHAPE.HEXAGON,
            'HORIZONTAL_SCROLL': MSO_SHAPE.HORIZONTAL_SCROLL,
            'ISOSCELES_TRIANGLE': MSO_SHAPE.ISOSCELES_TRIANGLE,
            
            # More Arrows
            'LEFT_ARROW': MSO_SHAPE.LEFT_ARROW,
            'LEFT_ARROW_CALLOUT': MSO_SHAPE.LEFT_ARROW_CALLOUT,
            'LEFT_BRACE': MSO_SHAPE.LEFT_BRACE,
            'LEFT_BRACKET': MSO_SHAPE.LEFT_BRACKET,
            'LEFT_CIRCULAR_ARROW': MSO_SHAPE.LEFT_CIRCULAR_ARROW,
            'LEFT_RIGHT_ARROW': MSO_SHAPE.LEFT_RIGHT_ARROW,
            'LEFT_RIGHT_ARROW_CALLOUT': MSO_SHAPE.LEFT_RIGHT_ARROW_CALLOUT,
            'LEFT_RIGHT_CIRCULAR_ARROW': MSO_SHAPE.LEFT_RIGHT_CIRCULAR_ARROW,
            'LEFT_RIGHT_RIBBON': MSO_SHAPE.LEFT_RIGHT_RIBBON,
            'LEFT_RIGHT_UP_ARROW': MSO_SHAPE.LEFT_RIGHT_UP_ARROW,
            'LEFT_UP_ARROW': MSO_SHAPE.LEFT_UP_ARROW,
            'LIGHTNING_BOLT': MSO_SHAPE.LIGHTNING_BOLT,
            
            # Line Callouts
            'LINE_CALLOUT_1': MSO_SHAPE.LINE_CALLOUT_1,
            'LINE_CALLOUT_1_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_1_ACCENT_BAR,
            'LINE_CALLOUT_1_BORDER_AND_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_1_BORDER_AND_ACCENT_BAR,
            'LINE_CALLOUT_1_NO_BORDER': MSO_SHAPE.LINE_CALLOUT_1_NO_BORDER,
            'LINE_CALLOUT_2': MSO_SHAPE.LINE_CALLOUT_2,
            'LINE_CALLOUT_2_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_2_ACCENT_BAR,
            'LINE_CALLOUT_2_BORDER_AND_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_2_BORDER_AND_ACCENT_BAR,
            'LINE_CALLOUT_2_NO_BORDER': MSO_SHAPE.LINE_CALLOUT_2_NO_BORDER,
            'LINE_CALLOUT_3': MSO_SHAPE.LINE_CALLOUT_3,
            'LINE_CALLOUT_3_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_3_ACCENT_BAR,
            'LINE_CALLOUT_3_BORDER_AND_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_3_BORDER_AND_ACCENT_BAR,
            'LINE_CALLOUT_3_NO_BORDER': MSO_SHAPE.LINE_CALLOUT_3_NO_BORDER,
            'LINE_CALLOUT_4': MSO_SHAPE.LINE_CALLOUT_4,
            'LINE_CALLOUT_4_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_4_ACCENT_BAR,
            'LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR': MSO_SHAPE.LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR,
            'LINE_CALLOUT_4_NO_BORDER': MSO_SHAPE.LINE_CALLOUT_4_NO_BORDER,
            'LINE_INVERSE': MSO_SHAPE.LINE_INVERSE,
            
            # Math Symbols
            'MATH_DIVIDE': MSO_SHAPE.MATH_DIVIDE,
            'MATH_EQUAL': MSO_SHAPE.MATH_EQUAL,
            'MATH_MINUS': MSO_SHAPE.MATH_MINUS,
            'MATH_MULTIPLY': MSO_SHAPE.MATH_MULTIPLY,
            'MATH_NOT_EQUAL': MSO_SHAPE.MATH_NOT_EQUAL,
            'MATH_PLUS': MSO_SHAPE.MATH_PLUS,
            
            # More Shapes
            'MOON': MSO_SHAPE.MOON,
            'NO_SYMBOL': MSO_SHAPE.NO_SYMBOL,
            'NON_ISOSCELES_TRAPEZOID': MSO_SHAPE.NON_ISOSCELES_TRAPEZOID,
            'NOTCHED_RIGHT_ARROW': MSO_SHAPE.NOTCHED_RIGHT_ARROW,
            'OCTAGON': MSO_SHAPE.OCTAGON,
            'OVAL': MSO_SHAPE.OVAL,
            'OVAL_CALLOUT': MSO_SHAPE.OVAL_CALLOUT,
            'PARALLELOGRAM': MSO_SHAPE.PARALLELOGRAM,
            'PENTAGON': MSO_SHAPE.PENTAGON,
            'PIE': MSO_SHAPE.PIE,
            'PIE_WEDGE': MSO_SHAPE.PIE_WEDGE,
            'PLAQUE': MSO_SHAPE.PLAQUE,
            'PLAQUE_TABS': MSO_SHAPE.PLAQUE_TABS,
            
            # Quad Arrows
            'QUAD_ARROW': MSO_SHAPE.QUAD_ARROW,
            'QUAD_ARROW_CALLOUT': MSO_SHAPE.QUAD_ARROW_CALLOUT,
            
            # Basic Rectangles
            'RECTANGLE': MSO_SHAPE.RECTANGLE,
            'RECTANGULAR_CALLOUT': MSO_SHAPE.RECTANGULAR_CALLOUT,
            'REGULAR_PENTAGON': MSO_SHAPE.REGULAR_PENTAGON,
            
            # Right Arrows
            'RIGHT_ARROW': MSO_SHAPE.RIGHT_ARROW,
            'RIGHT_ARROW_CALLOUT': MSO_SHAPE.RIGHT_ARROW_CALLOUT,
            'RIGHT_BRACE': MSO_SHAPE.RIGHT_BRACE,
            'RIGHT_BRACKET': MSO_SHAPE.RIGHT_BRACKET,
            'RIGHT_TRIANGLE': MSO_SHAPE.RIGHT_TRIANGLE,
            
            # Rounded Rectangles
            'ROUND_1_RECTANGLE': MSO_SHAPE.ROUND_1_RECTANGLE,
            'ROUND_2_DIAG_RECTANGLE': MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
            'ROUND_2_SAME_RECTANGLE': MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
            'ROUNDED_RECTANGLE': MSO_SHAPE.ROUNDED_RECTANGLE,
            'ROUNDED_RECTANGULAR_CALLOUT': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
            
            # Smiley and Snip Shapes
            'SMILEY_FACE': MSO_SHAPE.SMILEY_FACE,
            'SNIP_1_RECTANGLE': MSO_SHAPE.SNIP_1_RECTANGLE,
            'SNIP_2_DIAG_RECTANGLE': MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
            'SNIP_2_SAME_RECTANGLE': MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
            'SNIP_ROUND_RECTANGLE': MSO_SHAPE.SNIP_ROUND_RECTANGLE,
            'SQUARE_TABS': MSO_SHAPE.SQUARE_TABS,
            
            # Stars
            'STAR_10_POINT': MSO_SHAPE.STAR_10_POINT,
            'STAR_12_POINT': MSO_SHAPE.STAR_12_POINT,
            'STAR_16_POINT': MSO_SHAPE.STAR_16_POINT,
            'STAR_24_POINT': MSO_SHAPE.STAR_24_POINT,
            'STAR_32_POINT': MSO_SHAPE.STAR_32_POINT,
            'STAR_4_POINT': MSO_SHAPE.STAR_4_POINT,
            'STAR_5_POINT': MSO_SHAPE.STAR_5_POINT,
            'STAR_6_POINT': MSO_SHAPE.STAR_6_POINT,
            'STAR_7_POINT': MSO_SHAPE.STAR_7_POINT,
            'STAR_8_POINT': MSO_SHAPE.STAR_8_POINT,
            
            # More Arrows and Shapes
            'STRIPED_RIGHT_ARROW': MSO_SHAPE.STRIPED_RIGHT_ARROW,
            'SUN': MSO_SHAPE.SUN,
            'SWOOSH_ARROW': MSO_SHAPE.SWOOSH_ARROW,
            'TEAR': MSO_SHAPE.TEAR,
            'TRAPEZOID': MSO_SHAPE.TRAPEZOID,
            'U_TURN_ARROW': MSO_SHAPE.U_TURN_ARROW,
            
            # Up Arrows
            'UP_ARROW': MSO_SHAPE.UP_ARROW,
            'UP_ARROW_CALLOUT': MSO_SHAPE.UP_ARROW_CALLOUT,
            'UP_DOWN_ARROW': MSO_SHAPE.UP_DOWN_ARROW,
            'UP_DOWN_ARROW_CALLOUT': MSO_SHAPE.UP_DOWN_ARROW_CALLOUT,
            'UP_RIBBON': MSO_SHAPE.UP_RIBBON,
            
            # Scrolls and Wave
            'VERTICAL_SCROLL': MSO_SHAPE.VERTICAL_SCROLL,
            'WAVE': MSO_SHAPE.WAVE,
            
            # Chart Shapes (newer additions)
            'CHART_PLUS': MSO_SHAPE.CHART_PLUS,
            'CHART_STAR': MSO_SHAPE.CHART_STAR,
            'CHART_X': MSO_SHAPE.CHART_X,
            'DIAGONAL_STRIPE': MSO_SHAPE.DIAGONAL_STRIPE,
        }
        
        return auto_shape_mapping.get(shape_name, None)
    
    def map_shape_name_to_geometry(self, shape_name: str) -> str:
        """Map MSO_SHAPE names back to internal PowerPoint geometry names"""
        # Reverse mapping from MSO_SHAPE names to internal geometry names
        shape_to_geometry_mapping = {
            'RECTANGLE': 'rect',
            'ROUNDED_RECTANGLE': 'roundRect',
            'OVAL': 'ellipse',
            'ISOSCELES_TRIANGLE': 'triangle',
            'RIGHT_TRIANGLE': 'rtTriangle',
            'PARALLELOGRAM': 'parallelogram',
            'TRAPEZOID': 'trapezoid',
            'DIAMOND': 'diamond',
            'PENTAGON': 'pentagon',
            'HEXAGON': 'hexagon',
            'HEPTAGON': 'heptagon',
            'OCTAGON': 'octagon',
            'DECAGON': 'decagon',
            'DODECAGON': 'dodecagon',
            'PIE': 'pie',
            'CHORD': 'chord',
            'TEAR': 'teardrop',
            'FRAME': 'frame',
            'HALF_FRAME': 'halfFrame',
            'CORNER': 'corner',
            'DIAGONAL_STRIPE': 'diagStripe',
            'CROSS': 'plus',
            'PLAQUE': 'plaque',
            'CAN': 'can',
            'CUBE': 'cube',
            'BEVEL': 'bevel',
            'DONUT': 'donut',
            'NO_SYMBOL': 'noSmoking',
            'BLOCK_ARC': 'blockArc',
            'FOLDED_CORNER': 'foldedCorner',
            'SMILEY_FACE': 'smileyFace',
            'HEART': 'heart',
            'LIGHTNING_BOLT': 'lightningBolt',
            'SUN': 'sun',
            'MOON': 'moon',
            'ARC': 'arc',
            'DOUBLE_BRACKET': 'bracketPair',
            'DOUBLE_BRACE': 'bracePair',
            'LEFT_BRACKET': 'leftBracket',
            'RIGHT_BRACKET': 'rightBracket',
            'LEFT_BRACE': 'leftBrace',
            'RIGHT_BRACE': 'rightBrace',
            'RIGHT_ARROW': 'rightArrow',
            'LEFT_ARROW': 'leftArrow',
            'UP_ARROW': 'upArrow',
            'DOWN_ARROW': 'downArrow',
            'LEFT_RIGHT_ARROW': 'leftRightArrow',
            'UP_DOWN_ARROW': 'upDownArrow',
            'QUAD_ARROW': 'quadArrow',
            'LEFT_RIGHT_UP_ARROW': 'leftRightUpArrow',
            'BENT_ARROW': 'bentArrow',
            'U_TURN_ARROW': 'uturnArrow',
            'LEFT_UP_ARROW': 'leftUpArrow',
            'BENT_UP_ARROW': 'bentUpArrow',
            'CURVED_RIGHT_ARROW': 'curvedRightArrow',
            'CURVED_LEFT_ARROW': 'curvedLeftArrow',
            'CURVED_UP_ARROW': 'curvedUpArrow',
            'CURVED_DOWN_ARROW': 'curvedDownArrow',
            'STRIPED_RIGHT_ARROW': 'stripedRightArrow',
            'NOTCHED_RIGHT_ARROW': 'notchedRightArrow',
            'CIRCULAR_ARROW': 'circularArrow',
            'SWOOSH_ARROW': 'swooshArrow',
            'STAR_4_POINT': 'star4',
            'STAR_5_POINT': 'star5',
            'STAR_6_POINT': 'star6',
            'STAR_7_POINT': 'star7',
            'STAR_8_POINT': 'star8',
            'STAR_10_POINT': 'star10',
            'STAR_12_POINT': 'star12',
            'STAR_16_POINT': 'star16',
            'STAR_24_POINT': 'star24',
            'STAR_32_POINT': 'star32',
            'UP_RIBBON': 'ribbon2',
            'DOWN_RIBBON': 'ribbon',
            'CURVED_UP_RIBBON': 'ellipseRibbon2',
            'CURVED_DOWN_RIBBON': 'ellipseRibbon',
            'VERTICAL_SCROLL': 'verticalScroll',
            'HORIZONTAL_SCROLL': 'horizontalScroll',
            'WAVE': 'wave',
            'DOUBLE_WAVE': 'doubleWave',
            'GEAR_6': 'gear6',
            'GEAR_9': 'gear9',
            'FUNNEL': 'funnel',
            'MATH_PLUS': 'mathPlus',
            'MATH_MINUS': 'mathMinus',
            'MATH_MULTIPLY': 'mathMultiply',
            'MATH_DIVIDE': 'mathDivide',
            'MATH_EQUAL': 'mathEqual',
            'MATH_NOT_EQUAL': 'mathNotEqual',
            'CHEVRON': 'chevron',
            'PIE_WEDGE': 'pieWedge',
            'EXPLOSION1': 'irregularSeal1',
            'EXPLOSION2': 'irregularSeal2',
            'CLOUD': 'cloud',
            'CLOUD_CALLOUT': 'cloudCallout',
            'RECTANGULAR_CALLOUT': 'callout1',
            'ROUNDED_RECTANGULAR_CALLOUT': 'callout2',
            'OVAL_CALLOUT': 'callout3',
            'LINE_CALLOUT_1': 'borderCallout1',
            'LINE_CALLOUT_2': 'borderCallout2',
            'LINE_CALLOUT_3': 'borderCallout3',
            'LINE_CALLOUT_1_ACCENT_BAR': 'accentCallout1',
            'LINE_CALLOUT_2_ACCENT_BAR': 'accentCallout2',
            'LINE_CALLOUT_3_ACCENT_BAR': 'accentCallout3',
            'LINE_CALLOUT_4': 'callout90',
            'LINE_CALLOUT_4_ACCENT_BAR': 'accentCallout90',
            'LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR': 'borderCallout90',
        }
        
        return shape_to_geometry_mapping.get(shape_name.upper(), shape_name.lower())
    
    def get_mso_shape_from_type_number(self, type_number: int):
        """Convert MSO_SHAPE_TYPE number to appropriate MSO_SHAPE for creation"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
        
        # Complete mapping of MSO_SHAPE_TYPE numbers to appropriate MSO_SHAPE constants
        # for shape creation. Not all shape types can be created as autoshapes.
        shape_type_mapping = {
            -2: MSO_SHAPE.RECTANGLE,          # MIXED -> Default to rectangle
            1: MSO_SHAPE.RECTANGLE,           # AUTO_SHAPE -> Rectangle (most common autoshape)
            2: MSO_SHAPE.ROUNDED_RECTANGLE,   # CALLOUT -> Rounded rectangle as fallback
            3: None,                          # CHART -> Use add_chart instead
            4: MSO_SHAPE.RECTANGLE,           # COMMENT -> Rectangle as fallback
            5: MSO_SHAPE.ROUNDED_RECTANGLE,   # FREEFORM -> Rounded rectangle as fallback
            6: None,                          # GROUP -> Cannot create directly
            7: None,                          # EMBEDDED_OLE_OBJECT -> Cannot create
            8: MSO_SHAPE.RECTANGLE,           # FORM_CONTROL -> Rectangle as fallback
            9: None,                          # LINE -> Use add_connector instead
            10: None,                         # LINKED_OLE_OBJECT -> Cannot create
            11: None,                         # LINKED_PICTURE -> Use add_picture
            12: None,                         # OLE_CONTROL_OBJECT -> Cannot create
            13: None,                         # PICTURE -> Use add_picture instead
            14: None,                         # PLACEHOLDER -> Use existing placeholders
            15: MSO_SHAPE.RECTANGLE,          # TEXT_EFFECT -> Rectangle as fallback
            16: None,                         # MEDIA -> Cannot create
            17: None,                         # TEXT_BOX -> Use add_textbox instead
            18: MSO_SHAPE.RECTANGLE,          # SCRIPT_ANCHOR -> Rectangle as fallback
            19: None,                         # TABLE -> Use add_table instead
            20: MSO_SHAPE.RECTANGLE,          # CANVAS -> Rectangle as fallback
            21: MSO_SHAPE.DIAMOND,            # DIAGRAM -> Diamond as fallback
            22: MSO_SHAPE.OVAL,               # INK -> Oval as fallback
            23: MSO_SHAPE.RECTANGLE,          # INK_COMMENT -> Rectangle as fallback
            24: None,                         # IGX_GRAPHIC -> SmartArt, cannot create
            26: None,                         # WEB_VIDEO -> Cannot create
            27: MSO_SHAPE.RECTANGLE,          # CONTENT_APP -> Rectangle as fallback
            28: MSO_SHAPE.RECTANGLE,          # GRAPHIC -> Rectangle as fallback
            29: None,                         # LINKED_GRAPHIC -> Cannot create
            30: MSO_SHAPE.CUBE,               # 3D_MODEL -> Cube as fallback
            31: None,                         # LINKED_3D_MODEL -> Cannot create
        }
        
        # Return mapped shape or default to rectangle
        return shape_type_mapping.get(type_number, MSO_SHAPE.RECTANGLE)
    
    def get_layout_by_name(self, layout_name: str) -> Optional[Any]:
        """Get slide layout by name, fallback to first available layout"""
        # Try to find matching layout in presentation
        for i, layout in enumerate(self.presentation.slide_layouts):
            if layout.name.upper() == layout_name.upper():
                return layout
        
        # Fallback to appropriate layout based on name
        layout_mapping = {
            'TITLE': 0,
            'SECTION_HEADER': 1,
            'TITLE_AND_BODY': 1,
            'TITLE_AND_TWO_COLUMNS': 2,
            'TITLE_ONLY': 5,
            'ONE_COLUMN_TEXT': 1,
            'MAIN_POINT': 1,
            'SECTION_TITLE_AND_DESCRIPTION': 1,
            'CAPTION_ONLY': 6,
            'BIG_NUMBER': 1,
            'BLANK': 6
        }
        
        layout_index = layout_mapping.get(layout_name.upper(), 0)
        if layout_index < len(self.presentation.slide_layouts):
            return self.presentation.slide_layouts[layout_index]
        
        # Ultimate fallback
        return self.presentation.slide_layouts[0]
    
    def add_text_to_shape(self, shape, text: str):
        """Add text to a shape if it has a text frame"""
        if hasattr(shape, 'text_frame') and shape.text_frame is not None:
            shape.text_frame.text = text
            # Apply basic formatting
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
    
    def create_chart(self, slide, shape_info: Dict[str, Any]):
        """Create a chart based on chart data"""
        left = Inches(self.emu_to_inches(shape_info['left']))
        top = Inches(self.emu_to_inches(shape_info['top']))
        width = Inches(self.emu_to_inches(shape_info['width']))
        height = Inches(self.emu_to_inches(shape_info['height']))
        
        chart_data_info = shape_info.get('chart_data', {})
        
        try:
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            # Default chart type
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
            
            # Map chart type string to enum
            chart_type_str = chart_data_info.get('chart_type', '')
            if 'BAR' in chart_type_str.upper():
                chart_type = XL_CHART_TYPE.BAR_CLUSTERED
            elif 'LINE' in chart_type_str.upper():
                chart_type = XL_CHART_TYPE.LINE
            elif 'PIE' in chart_type_str.upper():
                chart_type = XL_CHART_TYPE.PIE
            elif 'AREA' in chart_type_str.upper():
                chart_type = XL_CHART_TYPE.AREA
            
            # Create chart data
            chart_data = CategoryChartData()
            
            # Add categories
            categories = chart_data_info.get('categories', ['Category 1', 'Category 2', 'Category 3'])
            chart_data.categories = categories
            
            # Add series
            series_list = chart_data_info.get('series', [])
            if not series_list:
                # Default series if no data
                chart_data.add_series('Series 1', [1, 2, 3])
            else:
                for series in series_list:
                    series_name = series.get('name', 'Series')
                    series_values = series.get('values', [1, 2, 3])
                    chart_data.add_series(series_name, series_values)
            
            # Add chart to slide
            chart = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
            
            # Set chart title if available
            title = chart_data_info.get('title')
            if title and chart.chart.has_title:
                chart.chart.chart_title.text_frame.text = title
            
            return chart.chart
            
        except Exception as e:
            print(f"Warning: Could not create chart: {str(e)}")
            # Fallback to text box with chart info
            return self.create_text_box(slide, shape_info)
    
    def create_table(self, slide, shape_info: Dict[str, Any]):
        """Create a table based on table data"""
        left = Inches(self.emu_to_inches(shape_info['left']))
        top = Inches(self.emu_to_inches(shape_info['top']))
        width = Inches(self.emu_to_inches(shape_info['width']))
        height = Inches(self.emu_to_inches(shape_info['height']))
        
        table_data_info = shape_info.get('table_data', {})
        
        try:
            # Get table dimensions
            rows = table_data_info.get('rows', 2)
            cols = table_data_info.get('columns', 2)
            data = table_data_info.get('data', [])
            
            # Ensure minimum table size
            if rows < 1:
                rows = 2
            if cols < 1:
                cols = 2
            
            # Create table
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # Fill table with data
            for row_idx in range(rows):
                for col_idx in range(cols):
                    cell = table.cell(row_idx, col_idx)
                    
                    # Get cell text from data if available
                    if row_idx < len(data) and col_idx < len(data[row_idx]):
                        cell_text = str(data[row_idx][col_idx])
                    else:
                        # Default cell content
                        if row_idx == 0:
                            cell_text = f"Header {col_idx + 1}"
                        else:
                            cell_text = f"Cell {row_idx},{col_idx + 1}"
                    
                    cell.text = cell_text
                    
                    # Format header row
                    if row_idx == 0:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.bold = True
                                run.font.size = Pt(12)
                    else:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(10)
            
            return table
            
        except Exception as e:
            print(f"Warning: Could not create table: {str(e)}")
            # Fallback to text box with table info
            return self.create_text_box(slide, shape_info)
    
    def create_text_box(self, slide, shape_info: Dict[str, Any]):
        """Create a text box shape"""
        left = Inches(self.emu_to_inches(shape_info['left']))
        top = Inches(self.emu_to_inches(shape_info['top']))
        width = Inches(self.emu_to_inches(shape_info['width']))
        height = Inches(self.emu_to_inches(shape_info['height']))
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = shape_info.get('text', '')
        
        # Apply formatting
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
        
        # Apply fill and line properties
        self.apply_fill_properties(text_box, shape_info.get('fill', {}))
        self.apply_line_properties(text_box, shape_info.get('line', {}))
        
        return text_box
    
    def create_shape_from_type(self, slide, shape_info: Dict[str, Any]):
        """Create a shape based on the shape type from JSON"""
        left = Inches(self.emu_to_inches(shape_info['left']))
        top = Inches(self.emu_to_inches(shape_info['top']))
        width = Inches(self.emu_to_inches(shape_info['width']))
        height = Inches(self.emu_to_inches(shape_info['height']))
        
        shape_type_str = shape_info.get('shape_type', '')
        mso_shape = None
        
        # First, try to parse shape type number and map to MSO_SHAPE
        type_number = self.parse_shape_type_number(shape_type_str)
        if type_number:
            mso_shape = self.get_mso_shape_from_type_number(type_number)
        
        # If no shape found by type number, try auto shape name mapping
        if not mso_shape:
            # Extract shape name without parentheses and numbers (handle duplicates)
            import re
            clean_name = re.sub(r'\s*\(\d+\)(\s*\(\d+\))*', '', shape_type_str).strip()
            mso_shape = self.get_mso_auto_shape_from_name(clean_name)
        
        # Create the shape
        if mso_shape:
            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
        else:
            # Fallback to rectangle for unsupported types
            from pptx.enum.shapes import MSO_SHAPE
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        
        # Apply fill properties
        self.apply_fill_properties(shape, shape_info.get('fill', {}))
        
        # Apply line properties
        self.apply_line_properties(shape, shape_info.get('line', {}))
        
        return shape
    
    def create_rectangle(self, slide, shape_info: Dict[str, Any]):
        """Create a rectangle shape (legacy method)"""
        return self.create_shape_from_type(slide, shape_info)
    
    def create_freeform_shape(self, slide, shape_info: Dict[str, Any]):
        """Create a freeform shape using proper shape type mapping"""
        # Use the new shape type mapping system
        return self.create_shape_from_type(slide, shape_info)
    
    def apply_fill_properties(self, shape, fill_info: Dict[str, Any]):
        """Apply fill properties to a shape"""
        if not fill_info:
            return
            
        try:
            from pptx.enum.dml import MSO_FILL_TYPE
            
            fill_type = fill_info.get('type', '')
            
            if 'SOLID' in fill_type or fill_info.get('solid'):
                shape.fill.solid()
                fore_color = fill_info.get('fore_color', {})
                self.apply_color_properties(shape.fill.fore_color, fore_color)
                
            elif 'PATTERN' in fill_type or fill_info.get('pattern'):
                shape.fill.patterned()
                fore_color = fill_info.get('fore_color', {})
                back_color = fill_info.get('back_color', {})
                self.apply_color_properties(shape.fill.fore_color, fore_color)
                self.apply_color_properties(shape.fill.back_color, back_color)
                
            elif 'BACKGROUND' in fill_type or fill_info.get('background'):
                shape.fill.background()
                
        except Exception as e:
            print(f"Warning: Could not apply fill properties: {str(e)}")
    
    def apply_line_properties(self, shape, line_info: Dict[str, Any]):
        """Apply line/border properties to a shape"""
        if not line_info:
            return
            
        try:
            line = shape.line
            
            # Apply line width
            if 'width' in line_info and line_info['width'] is not None:
                from pptx.util import Emu
                line.width = Emu(line_info['width'])
            
            # Apply line color
            color_info = line_info.get('color', {})
            if color_info:
                self.apply_color_properties(line.color, color_info)
                
        except Exception as e:
            print(f"Warning: Could not apply line properties: {str(e)}")
    
    def apply_color_properties(self, color_obj, color_info: Dict[str, Any]):
        """Apply color properties to a color object"""
        if not color_info:
            return
            
        try:
            from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
            
            # Apply RGB color
            rgb_info = color_info.get('rgb', {})
            if rgb_info:
                # Handle hex color if RGB values are null
                if rgb_info.get('hex') and (rgb_info.get('red') is None or rgb_info.get('green') is None or rgb_info.get('blue') is None):
                    hex_color = rgb_info.get('hex', '000000')
                    # Remove # if present
                    hex_color = hex_color.replace('#', '')
                    if len(hex_color) == 6:
                        red = int(hex_color[0:2], 16)
                        green = int(hex_color[2:4], 16)
                        blue = int(hex_color[4:6], 16)
                        color_obj.rgb = RGBColor(red, green, blue)
                elif rgb_info.get('red') is not None:
                    red = rgb_info.get('red', 0)
                    green = rgb_info.get('green', 0)
                    blue = rgb_info.get('blue', 0)
                    color_obj.rgb = RGBColor(red, green, blue)
                
            # Apply theme color
            elif 'theme_color' in color_info:
                theme_color_str = color_info['theme_color']
                # Map theme color string to enum (simplified mapping)
                theme_mapping = {
                    'ACCENT_1': MSO_THEME_COLOR.ACCENT_1,
                    'ACCENT_2': MSO_THEME_COLOR.ACCENT_2,
                    'ACCENT_3': MSO_THEME_COLOR.ACCENT_3,
                    'ACCENT_4': MSO_THEME_COLOR.ACCENT_4,
                    'ACCENT_5': MSO_THEME_COLOR.ACCENT_5,
                    'ACCENT_6': MSO_THEME_COLOR.ACCENT_6,
                    'BACKGROUND_1': MSO_THEME_COLOR.BACKGROUND_1,
                    'BACKGROUND_2': MSO_THEME_COLOR.BACKGROUND_2,
                    'DARK_1': MSO_THEME_COLOR.DARK_1,
                    'DARK_2': MSO_THEME_COLOR.DARK_2,
                    'FOLLOWED_HYPERLINK': MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
                    'HYPERLINK': MSO_THEME_COLOR.HYPERLINK,
                    'LIGHT_1': MSO_THEME_COLOR.LIGHT_1,
                    'LIGHT_2': MSO_THEME_COLOR.LIGHT_2,
                    'TEXT_1': MSO_THEME_COLOR.TEXT_1,
                    'TEXT_2': MSO_THEME_COLOR.TEXT_2,
                }
                
                for key, value in theme_mapping.items():
                    if key in theme_color_str:
                        color_obj.theme_color = value
                        break
            
            # Apply brightness adjustment
            if 'brightness' in color_info and color_info['brightness'] is not None:
                color_obj.brightness = color_info['brightness']
                
        except Exception as e:
            print(f"Warning: Could not apply color properties: {str(e)}")
    
    def apply_theme_to_presentation(self):
        """Apply theme settings to the presentation"""
        # Since theme extraction had errors, we'll apply basic styling
        print("Applying basic theme styling...")
        
        # Set slide size (standard 16:9)
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(7.5)
    
    def generate_slides(self):
        """Generate slides based on shapes data"""
        self.apply_theme_to_presentation()
        
        for slide_data in self.shapes_data:
            slide_index = slide_data['slide_index']
            shapes = slide_data['shapes']
            
            print(f"Creating slide {slide_index + 1} with {len(shapes)} shapes...")
            
            # Determine the best layout based on shapes
            layout_name = 'BLANK'  # Default to blank layout
            if shapes:
                # Simple heuristic: if we have placeholder shapes, try to use appropriate layout
                placeholder_count = sum(1 for s in shapes if 'PLACEHOLDER' in s.get('shape_type', ''))
                if placeholder_count >= 2:
                    layout_name = 'TITLE_AND_BODY'
                elif placeholder_count == 1:
                    layout_name = 'TITLE_ONLY'
            
            # Get the layout
            layout = self.get_layout_by_name(layout_name)
            slide = self.presentation.slides.add_slide(layout)
            
            # Add shapes to the slide
            for shape_info in shapes:
                self.add_shape_to_slide(slide, shape_info)
    
    def add_shape_to_slide(self, slide, shape_info: Dict[str, Any]):
        """Add a shape to the slide based on shape info"""
        shape_type = shape_info.get('shape_type', '')
        text = shape_info.get('text', '')
        
        try:
            # Check for special shape types that need different handling
            if 'CHART' in shape_type:
                # Create chart from chart data
                self.create_chart(slide, shape_info)
            elif 'TABLE' in shape_type:
                # Create table from table data
                self.create_table(slide, shape_info)
            elif 'PLACEHOLDER' in shape_type:
                # Handle placeholder with special content types
                self.handle_placeholder(slide, shape_info)
            elif 'TEXT_BOX' in shape_type or (text and 'AUTO_SHAPE' not in shape_type):
                # Create text box for text-focused shapes
                self.create_text_box(slide, shape_info)
            else:
                # Create appropriate shape based on type mapping
                shape = self.create_shape_from_type(slide, shape_info)
                if text:
                    self.add_text_to_shape(shape, text)
        except Exception as e:
            print(f"Warning: Could not create shape {shape_info.get('name', 'Unknown')}: {str(e)}")
    
    def handle_placeholder(self, slide, shape_info: Dict[str, Any]):
        """Handle placeholder shapes with special content types"""
        placeholder_info = shape_info.get('placeholder_info', {})
        
        # Check if placeholder has chart or table data
        if 'chart_data' in shape_info:
            # Try to insert chart into placeholder
            self.insert_chart_into_placeholder(slide, shape_info)
        elif 'table_data' in shape_info:
            # Try to insert table into placeholder
            self.insert_table_into_placeholder(slide, shape_info)
        else:
            # Regular text placeholder
            self.fill_placeholder(slide, shape_info)
    
    def insert_chart_into_placeholder(self, slide, shape_info: Dict[str, Any]):
        """Insert chart into a placeholder if possible"""
        placeholder_info = shape_info.get('placeholder_info', {})
        placeholder_type = placeholder_info.get('placeholder_type', '')
        
        # Try to find appropriate placeholder for chart
        for placeholder in slide.placeholders:
            try:
                if hasattr(placeholder, 'insert_chart'):
                    # Use the chart creation method
                    chart_data_info = shape_info.get('chart_data', {})
                    
                    from pptx.chart.data import CategoryChartData
                    from pptx.enum.chart import XL_CHART_TYPE
                    
                    chart_data = CategoryChartData()
                    categories = chart_data_info.get('categories', ['Cat 1', 'Cat 2', 'Cat 3'])
                    chart_data.categories = categories
                    
                    series_list = chart_data_info.get('series', [])
                    if not series_list:
                        chart_data.add_series('Series 1', [1, 2, 3])
                    else:
                        for series in series_list:
                            chart_data.add_series(series.get('name', 'Series'), series.get('values', [1, 2, 3]))
                    
                    placeholder.insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
                    return
            except:
                continue
        
        # Fallback: create chart as regular shape
        self.create_chart(slide, shape_info)
    
    def insert_table_into_placeholder(self, slide, shape_info: Dict[str, Any]):
        """Insert table into a placeholder if possible"""
        table_data_info = shape_info.get('table_data', {})
        
        # Try to find appropriate placeholder for table
        for placeholder in slide.placeholders:
            try:
                if hasattr(placeholder, 'insert_table'):
                    rows = table_data_info.get('rows', 2)
                    cols = table_data_info.get('columns', 2)
                    data = table_data_info.get('data', [])
                    
                    table = placeholder.insert_table(rows, cols)
                    
                    # Fill table with data
                    for row_idx in range(min(rows, len(data))):
                        for col_idx in range(min(cols, len(data[row_idx]))):
                            cell = table.table.cell(row_idx, col_idx)
                            cell.text = str(data[row_idx][col_idx])
                    
                    return
            except:
                continue
        
        # Fallback: create table as regular shape
        self.create_table(slide, shape_info)
    
    def fill_placeholder(self, slide, shape_info: Dict[str, Any]):
        """Try to fill slide placeholders with content"""
        text = shape_info.get('text', '')
        if not text:
            return
        
        # Try to find appropriate placeholder
        for placeholder in slide.placeholders:
            if hasattr(placeholder, 'text_frame') and placeholder.text_frame is not None:
                try:
                    if not placeholder.text_frame.text:  # Only fill empty placeholders
                        placeholder.text_frame.text = text
                        # Apply formatting
                        for paragraph in placeholder.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(18)
                        return
                except:
                    continue
        
        # If no placeholder available, create a text box
        self.create_text_box(slide, shape_info)
    
    def save_presentation(self, output_file: str):
        """Save the presentation to file"""
        self.presentation.save(output_file)
        print(f"Presentation saved to: {output_file}")


def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint presentation from JSON files')
    parser.add_argument('shapes_file', help='Path to shapes JSON file')
    parser.add_argument('layouts_file', help='Path to layouts JSON file')
    parser.add_argument('theme_file', help='Path to theme JSON file')
    parser.add_argument('--output', '-o', default='generated_presentation.pptx', 
                       help='Output PowerPoint file name (default: generated_presentation.pptx)')
    
    args = parser.parse_args()
    
    # Validate input files
    for file_path in [args.shapes_file, args.layouts_file, args.theme_file]:
        if not Path(file_path).exists():
            print(f"Error: File '{file_path}' does not exist.")
            sys.exit(1)
    
    try:
        # Initialize generator
        generator = PPTGenerator()
        
        # Load JSON data
        print("Loading JSON files...")
        generator.load_json_files(args.shapes_file, args.layouts_file, args.theme_file)
        
        # Generate slides
        print("Generating presentation...")
        generator.generate_slides()
        
        # Save presentation
        generator.save_presentation(args.output)
        
        print(f"\nPresentation generation completed successfully!")
        print(f"Output file: {args.output}")
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()