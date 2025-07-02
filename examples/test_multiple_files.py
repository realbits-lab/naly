#!/usr/bin/env python3

import os
import sys
import json
import time
import argparse
import subprocess
from pathlib import Path
from typing import Dict, List, Tuple, Any
from datetime import datetime
import traceback


class MultiFileValidator:
    """Validate PowerPoint extraction and generation across multiple files"""
    
    def __init__(self, test_dir: str = "test_presentations"):
        self.test_dir = Path(test_dir)
        self.results = {
            'test_date': datetime.now().isoformat(),
            'test_files': [],
            'summary': {
                'total_files': 0,
                'successful': 0,
                'failed': 0,
                'average_fidelity': 0.0,
                'extraction_time': 0.0,
                'generation_time': 0.0,
            },
            'detailed_results': []
        }
        
    def create_test_presentations(self):
        """Create various test PowerPoint files for validation"""
        print("Creating test presentations...")
        
        # Ensure test directory exists
        self.test_dir.mkdir(exist_ok=True)
        
        test_configs = [
            {
                'name': 'simple_shapes.pptx',
                'description': 'Basic shapes and text',
                'slides': [
                    {
                        'layout': 'Title Slide',
                        'content': {
                            'title': 'Simple Shapes Test',
                            'subtitle': 'Testing basic shape rendering'
                        }
                    },
                    {
                        'layout': 'Blank',
                        'shapes': ['rectangle', 'oval', 'triangle', 'arrow']
                    }
                ]
            },
            {
                'name': 'charts_tables.pptx',
                'description': 'Charts and tables',
                'slides': [
                    {
                        'layout': 'Title and Content',
                        'content': {
                            'title': 'Data Visualization Test',
                            'chart': 'column_chart'
                        }
                    },
                    {
                        'layout': 'Title and Content',
                        'content': {
                            'title': 'Table Test',
                            'table': {'rows': 3, 'cols': 3}
                        }
                    }
                ]
            },
            {
                'name': 'themed_presentation.pptx',
                'description': 'Theme colors and fonts',
                'slides': [
                    {
                        'layout': 'Title Slide',
                        'content': {
                            'title': 'Theme Test',
                            'subtitle': 'Testing theme preservation'
                        },
                        'theme_colors': True
                    }
                ]
            },
            {
                'name': 'complex_layouts.pptx',
                'description': 'Multiple layout types',
                'slides': [
                    {
                        'layout': 'Title Slide',
                        'content': {'title': 'Layout Test'}
                    },
                    {
                        'layout': 'Title and Content',
                        'content': {'title': 'Content Layout'}
                    },
                    {
                        'layout': 'Section Header',
                        'content': {'title': 'Section'}
                    },
                    {
                        'layout': 'Two Content',
                        'content': {'title': 'Two Columns'}
                    }
                ]
            }
        ]
        
        # Generate test presentations using python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            for config in test_configs:
                prs = Presentation()
                
                for slide_config in config['slides']:
                    # Get layout
                    layout_name = slide_config['layout']
                    layout = None
                    for idx, slide_layout in enumerate(prs.slide_layouts):
                        if layout_name.lower() in slide_layout.name.lower():
                            layout = slide_layout
                            break
                    
                    if not layout:
                        layout = prs.slide_layouts[0]  # Default to first layout
                    
                    slide = prs.slides.add_slide(layout)
                    
                    # Add content based on configuration
                    if 'content' in slide_config:
                        content = slide_config['content']
                        
                        # Add title
                        if 'title' in content and slide.shapes.title:
                            slide.shapes.title.text = content['title']
                        
                        # Add subtitle
                        if 'subtitle' in content:
                            for shape in slide.placeholders:
                                if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                                    shape.text = content['subtitle']
                        
                        # Add chart
                        if 'chart' in content:
                            chart_data = CategoryChartData()
                            chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
                            chart_data.add_series('Series 1', (19.2, 21.4, 16.7, 22.3))
                            
                            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
                            slide.shapes.add_chart(
                                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
                            )
                        
                        # Add table
                        if 'table' in content:
                            table_info = content['table']
                            rows = table_info.get('rows', 2)
                            cols = table_info.get('cols', 2)
                            
                            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
                            table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table
                            
                            # Fill with sample data
                            for row in range(rows):
                                for col in range(cols):
                                    cell = table.cell(row, col)
                                    if row == 0:
                                        cell.text = f"Header {col + 1}"
                                    else:
                                        cell.text = f"Data {row},{col + 1}"
                    
                    # Add shapes
                    if 'shapes' in slide_config:
                        shape_positions = [
                            (Inches(1), Inches(2)),
                            (Inches(4), Inches(2)),
                            (Inches(1), Inches(4)),
                            (Inches(4), Inches(4))
                        ]
                        
                        shape_mapping = {
                            'rectangle': MSO_SHAPE.RECTANGLE,
                            'oval': MSO_SHAPE.OVAL,
                            'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                            'arrow': MSO_SHAPE.RIGHT_ARROW
                        }
                        
                        for idx, shape_name in enumerate(slide_config['shapes'][:4]):
                            if shape_name in shape_mapping and idx < len(shape_positions):
                                x, y = shape_positions[idx]
                                shape = slide.shapes.add_shape(
                                    shape_mapping[shape_name],
                                    x, y, Inches(2), Inches(1.5)
                                )
                                shape.text = shape_name.capitalize()
                
                # Save presentation
                output_path = self.test_dir / config['name']
                prs.save(str(output_path))
                print(f"  Created: {config['name']} - {config['description']}")
                
        except ImportError:
            print("Warning: python-pptx not installed. Using existing test files.")
        except Exception as e:
            print(f"Error creating test presentations: {str(e)}")
    
    def validate_file(self, pptx_file: Path, use_enhanced: bool = True) -> Dict[str, Any]:
        """Validate a single PowerPoint file through extraction and regeneration"""
        result = {
            'file': str(pptx_file),
            'status': 'pending',
            'extraction_time': 0,
            'generation_time': 0,
            'comparison_time': 0,
            'fidelity_score': 0,
            'errors': []
        }
        
        try:
            # Create output directory for this file
            output_dir = self.test_dir / f"output_{pptx_file.stem}"
            output_dir.mkdir(exist_ok=True)
            
            # Step 1: Extract data
            print(f"\n  Extracting {pptx_file.name}...")
            start_time = time.time()
            
            extract_cmd = [
                sys.executable, 'ppt_extractor.py',
                str(pptx_file),
                '--output-dir', str(output_dir)
            ]
            
            extract_result = subprocess.run(
                extract_cmd,
                capture_output=True,
                text=True
            )
            
            extraction_time = time.time() - start_time
            result['extraction_time'] = extraction_time
            
            if extract_result.returncode != 0:
                result['status'] = 'extraction_failed'
                result['errors'].append(f"Extraction error: {extract_result.stderr}")
                return result
            
            # Step 2: Generate PowerPoint
            print(f"  Generating from extracted data...")
            start_time = time.time()
            
            generator_script = 'ppt_generator_enhanced.py' if use_enhanced else 'ppt_generator.py'
            shapes_file = output_dir / f"{pptx_file.stem}_shapes.json"
            layouts_file = output_dir / f"{pptx_file.stem}_layouts.json"
            theme_file = output_dir / f"{pptx_file.stem}_theme.json"
            output_file = output_dir / f"regenerated_{pptx_file.name}"
            
            generate_cmd = [
                sys.executable, generator_script,
                str(shapes_file),
                str(layouts_file),
                str(theme_file),
                '--output', str(output_file)
            ]
            
            generate_result = subprocess.run(
                generate_cmd,
                capture_output=True,
                text=True
            )
            
            generation_time = time.time() - start_time
            result['generation_time'] = generation_time
            
            if generate_result.returncode != 0:
                result['status'] = 'generation_failed'
                result['errors'].append(f"Generation error: {generate_result.stderr}")
                return result
            
            # Step 3: Compare structures
            print(f"  Comparing structures...")
            start_time = time.time()
            
            comparison_output = output_dir / 'comparison_report.json'
            compare_cmd = [
                sys.executable, 'structure_comparison.py',
                str(pptx_file),
                str(output_file),
                '--output', str(comparison_output)
            ]
            
            compare_result = subprocess.run(
                compare_cmd,
                capture_output=True,
                text=True
            )
            
            comparison_time = time.time() - start_time
            result['comparison_time'] = comparison_time
            
            if compare_result.returncode != 0:
                result['status'] = 'comparison_failed'
                result['errors'].append(f"Comparison error: {compare_result.stderr}")
                return result
            
            # Load comparison results
            if comparison_output.exists():
                with open(comparison_output, 'r') as f:
                    comparison_data = json.load(f)
                    result['fidelity_score'] = comparison_data.get('fidelity_score', 0)
                    result['comparison_summary'] = comparison_data.get('summary', {})
            
            result['status'] = 'success'
            
        except Exception as e:
            result['status'] = 'error'
            result['errors'].append(f"Unexpected error: {str(e)}")
            traceback.print_exc()
        
        return result
    
    def run_validation(self, test_files: List[str] = None, use_enhanced: bool = True):
        """Run validation on multiple files"""
        print("\n" + "="*80)
        print("MULTI-FILE POWERPOINT VALIDATION")
        print("="*80)
        
        # Get test files
        if test_files:
            pptx_files = [Path(f) for f in test_files if Path(f).exists()]
        else:
            # Find all PPTX files in test directory
            pptx_files = list(self.test_dir.glob('*.pptx'))
            
            # Also check current directory
            if not pptx_files:
                pptx_files = list(Path('.').glob('*.pptx'))
        
        if not pptx_files:
            print("No PowerPoint files found for testing.")
            print("Creating test presentations...")
            self.create_test_presentations()
            pptx_files = list(self.test_dir.glob('*.pptx'))
        
        self.results['test_files'] = [str(f) for f in pptx_files]
        self.results['summary']['total_files'] = len(pptx_files)
        
        print(f"\nTesting {len(pptx_files)} PowerPoint files...")
        print(f"Using {'enhanced' if use_enhanced else 'original'} generator\n")
        
        # Validate each file
        total_fidelity = 0
        total_extraction_time = 0
        total_generation_time = 0
        
        for pptx_file in pptx_files:
            print(f"Processing: {pptx_file.name}")
            result = self.validate_file(pptx_file, use_enhanced)
            
            self.results['detailed_results'].append(result)
            
            if result['status'] == 'success':
                self.results['summary']['successful'] += 1
                total_fidelity += result['fidelity_score']
                total_extraction_time += result['extraction_time']
                total_generation_time += result['generation_time']
                print(f"  ✓ Success - Fidelity: {result['fidelity_score']:.1%}")
            else:
                self.results['summary']['failed'] += 1
                print(f"  ✗ Failed - Status: {result['status']}")
        
        # Calculate summary statistics
        if self.results['summary']['successful'] > 0:
            self.results['summary']['average_fidelity'] = (
                total_fidelity / self.results['summary']['successful']
            )
            self.results['summary']['extraction_time'] = total_extraction_time
            self.results['summary']['generation_time'] = total_generation_time
        
        # Print summary
        self.print_summary()
        
        # Save results
        results_file = self.test_dir / f"validation_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        with open(results_file, 'w') as f:
            json.dump(self.results, f, indent=2)
        print(f"\nDetailed results saved to: {results_file}")
    
    def print_summary(self):
        """Print validation summary"""
        summary = self.results['summary']
        
        print("\n" + "="*80)
        print("VALIDATION SUMMARY")
        print("="*80)
        
        print(f"\nTotal Files Tested: {summary['total_files']}")
        print(f"Successful: {summary['successful']}")
        print(f"Failed: {summary['failed']}")
        
        if summary['successful'] > 0:
            print(f"\nAverage Fidelity Score: {summary['average_fidelity']:.1%}")
            print(f"Total Extraction Time: {summary['extraction_time']:.2f}s")
            print(f"Total Generation Time: {summary['generation_time']:.2f}s")
            print(f"Average Time per File: {(summary['extraction_time'] + summary['generation_time']) / summary['successful']:.2f}s")
        
        # Show individual file results
        print("\nIndividual File Results:")
        for result in self.results['detailed_results']:
            status_symbol = "✓" if result['status'] == 'success' else "✗"
            print(f"  {status_symbol} {Path(result['file']).name}: ", end="")
            
            if result['status'] == 'success':
                print(f"Fidelity {result['fidelity_score']:.1%} "
                      f"(Extract: {result['extraction_time']:.1f}s, "
                      f"Generate: {result['generation_time']:.1f}s)")
            else:
                print(f"{result['status']}")
                if result['errors']:
                    print(f"    Errors: {result['errors'][0][:100]}...")
        
        # Recommendations
        print("\nRecommendations:")
        if summary['average_fidelity'] < 0.5:
            print("  - Low average fidelity. Review extraction and generation logic.")
        elif summary['average_fidelity'] < 0.8:
            print("  - Moderate fidelity. Focus on specific feature improvements.")
        else:
            print("  - Good fidelity achieved. Fine-tune remaining issues.")
        
        failed_ratio = summary['failed'] / summary['total_files'] if summary['total_files'] > 0 else 0
        if failed_ratio > 0.2:
            print("  - High failure rate. Check for compatibility issues.")


def main():
    parser = argparse.ArgumentParser(
        description='Validate PowerPoint extraction and generation across multiple files'
    )
    parser.add_argument('files', nargs='*', help='PowerPoint files to test (optional)')
    parser.add_argument('--test-dir', default='test_presentations',
                       help='Directory for test files and output (default: test_presentations)')
    parser.add_argument('--create-tests', action='store_true',
                       help='Create test presentations')
    parser.add_argument('--use-original', action='store_true',
                       help='Use original generator instead of enhanced')
    
    args = parser.parse_args()
    
    # Create validator
    validator = MultiFileValidator(args.test_dir)
    
    # Create test presentations if requested
    if args.create_tests:
        validator.create_test_presentations()
        if not args.files:  # If no files specified, exit after creation
            print("\nTest presentations created. Run without --create-tests to validate them.")
            return
    
    # Run validation
    use_enhanced = not args.use_original
    validator.run_validation(args.files, use_enhanced)


if __name__ == "__main__":
    main()