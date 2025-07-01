#!/usr/bin/env python3

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide, SlideLayout
from typing import Dict, List, Any


class PPTExtractor:
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.presentation = Presentation(file_path)
        
    def extract_fill_properties(self, fill) -> Dict[str, Any]:
        """Extract detailed fill properties including colors"""
        fill_info = {
            'type': str(fill.type) if fill.type else None
        }
        
        try:
            from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE
            
            if fill.type == MSO_FILL_TYPE.SOLID:
                fill_info['solid'] = True
                color_info = self.extract_color_properties(fill.fore_color)
                fill_info['fore_color'] = color_info
                
            elif fill.type == MSO_FILL_TYPE.PATTERN:
                fill_info['pattern'] = True
                fill_info['fore_color'] = self.extract_color_properties(fill.fore_color)
                fill_info['back_color'] = self.extract_color_properties(fill.back_color)
                
            elif fill.type == MSO_FILL_TYPE.GRADIENT:
                fill_info['gradient'] = True
                
            elif fill.type == MSO_FILL_TYPE.PICTURE:
                fill_info['picture'] = True
                
            elif fill.type == MSO_FILL_TYPE.BACKGROUND:
                fill_info['background'] = True
                
        except Exception as e:
            fill_info['error'] = f"Could not extract fill properties: {str(e)}"
            
        return fill_info
    
    def extract_color_properties(self, color) -> Dict[str, Any]:
        """Extract color properties including RGB values"""
        color_info = {
            'type': str(color.type) if hasattr(color, 'type') and color.type else None
        }
        
        try:
            from pptx.enum.dml import MSO_COLOR_TYPE
            
            if hasattr(color, 'type') and color.type == MSO_COLOR_TYPE.RGB:
                if hasattr(color, 'rgb') and color.rgb:
                    rgb = color.rgb
                    color_info['rgb'] = {
                        'hex': str(rgb),
                        'red': rgb.red if hasattr(rgb, 'red') else None,
                        'green': rgb.green if hasattr(rgb, 'green') else None,
                        'blue': rgb.blue if hasattr(rgb, 'blue') else None
                    }
                    
            elif hasattr(color, 'type') and color.type == MSO_COLOR_TYPE.SCHEME:
                if hasattr(color, 'theme_color'):
                    color_info['theme_color'] = str(color.theme_color)
                    
            # Extract brightness if available
            if hasattr(color, 'brightness') and color.brightness is not None:
                color_info['brightness'] = color.brightness
                
        except Exception as e:
            color_info['error'] = f"Could not extract color properties: {str(e)}"
            
        return color_info
    
    def extract_line_properties(self, line) -> Dict[str, Any]:
        """Extract line/border properties"""
        line_info = {}
        
        try:
            # Line width
            if hasattr(line, 'width') and line.width is not None:
                line_info['width'] = line.width
                
            # Line color
            if hasattr(line, 'color'):
                color_info = self.extract_color_properties(line.color)
                line_info['color'] = color_info
                
            # Line fill (for compound lines)
            if hasattr(line, 'fill'):
                fill_info = self.extract_fill_properties(line.fill)
                line_info['fill'] = fill_info
                
        except Exception as e:
            line_info['error'] = f"Could not extract line properties: {str(e)}"
            
        return line_info

    def get_auto_shape_type(self, shape) -> str:
        """Extract the specific auto shape type for MSO_SHAPE_TYPE.AUTO_SHAPE"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Try to get the auto shape type from the shape element
                if hasattr(shape, '_element') and hasattr(shape._element, 'prstGeom'):
                    prst_geom = shape._element.prstGeom
                    if hasattr(prst_geom, 'prst') and prst_geom.prst:
                        # Map internal geometry names to MSO_SHAPE names
                        geom_name = prst_geom.prst
                        return self.map_geometry_to_shape_name(geom_name)
                
                # If we can't determine the specific type, try to get it from shape properties
                if hasattr(shape, 'auto_shape_type'):
                    return str(shape.auto_shape_type)
                    
                # Fallback: return the generic AUTO_SHAPE with number
                return f"AUTO_SHAPE ({shape.shape_type.value})"
            else:
                # For non-auto shapes, return the shape type with number (avoid duplication)
                shape_name = str(shape.shape_type).split('.')[-1]
                return f"{shape_name} ({shape.shape_type.value})"
                
        except Exception:
            # Fallback to original behavior
            return str(shape.shape_type)
    
    def map_geometry_to_shape_name(self, geom_name: str) -> str:
        """Map internal PowerPoint geometry names to MSO_SHAPE names"""
        # Mapping from internal geometry names to MSO_SHAPE constants
        geometry_mapping = {
            'rect': 'RECTANGLE',
            'roundRect': 'ROUNDED_RECTANGLE',
            'ellipse': 'OVAL',
            'triangle': 'ISOSCELES_TRIANGLE',
            'rtTriangle': 'RIGHT_TRIANGLE',
            'parallelogram': 'PARALLELOGRAM',
            'trapezoid': 'TRAPEZOID',
            'diamond': 'DIAMOND',
            'pentagon': 'PENTAGON',
            'hexagon': 'HEXAGON',
            'heptagon': 'HEPTAGON',
            'octagon': 'OCTAGON',
            'decagon': 'DECAGON',
            'dodecagon': 'DODECAGON',
            'pie': 'PIE',
            'chord': 'CHORD',
            'teardrop': 'TEAR',
            'frame': 'FRAME',
            'halfFrame': 'HALF_FRAME',
            'corner': 'CORNER',
            'diagStripe': 'DIAGONAL_STRIPE',
            'plus': 'CROSS',
            'plaque': 'PLAQUE',
            'can': 'CAN',
            'cube': 'CUBE',
            'bevel': 'BEVEL',
            'donut': 'DONUT',
            'noSmoking': 'NO_SYMBOL',
            'blockArc': 'BLOCK_ARC',
            'foldedCorner': 'FOLDED_CORNER',
            'smileyFace': 'SMILEY_FACE',
            'heart': 'HEART',
            'lightningBolt': 'LIGHTNING_BOLT',
            'sun': 'SUN',
            'moon': 'MOON',
            'arc': 'ARC',
            'bracketPair': 'DOUBLE_BRACKET',
            'bracePair': 'DOUBLE_BRACE',
            'leftBracket': 'LEFT_BRACKET',
            'rightBracket': 'RIGHT_BRACKET',
            'leftBrace': 'LEFT_BRACE',
            'rightBrace': 'RIGHT_BRACE',
            'rightArrow': 'RIGHT_ARROW',
            'leftArrow': 'LEFT_ARROW',
            'upArrow': 'UP_ARROW',
            'downArrow': 'DOWN_ARROW',
            'leftRightArrow': 'LEFT_RIGHT_ARROW',
            'upDownArrow': 'UP_DOWN_ARROW',
            'quadArrow': 'QUAD_ARROW',
            'leftRightUpArrow': 'LEFT_RIGHT_UP_ARROW',
            'bentArrow': 'BENT_ARROW',
            'uturnArrow': 'U_TURN_ARROW',
            'leftUpArrow': 'LEFT_UP_ARROW',
            'bentUpArrow': 'BENT_UP_ARROW',
            'curvedRightArrow': 'CURVED_RIGHT_ARROW',
            'curvedLeftArrow': 'CURVED_LEFT_ARROW',
            'curvedUpArrow': 'CURVED_UP_ARROW',
            'curvedDownArrow': 'CURVED_DOWN_ARROW',
            'stripedRightArrow': 'STRIPED_RIGHT_ARROW',
            'notchedRightArrow': 'NOTCHED_RIGHT_ARROW',
            'circularArrow': 'CIRCULAR_ARROW',
            'swooshArrow': 'SWOOSH_ARROW',
            'star4': 'STAR_4_POINT',
            'star5': 'STAR_5_POINT',
            'star6': 'STAR_6_POINT',
            'star7': 'STAR_7_POINT',
            'star8': 'STAR_8_POINT',
            'star10': 'STAR_10_POINT',
            'star12': 'STAR_12_POINT',
            'star16': 'STAR_16_POINT',
            'star24': 'STAR_24_POINT',
            'star32': 'STAR_32_POINT',
            'ribbon2': 'UP_RIBBON',
            'ribbon': 'DOWN_RIBBON',
            'ellipseRibbon2': 'CURVED_UP_RIBBON',
            'ellipseRibbon': 'CURVED_DOWN_RIBBON',
            'verticalScroll': 'VERTICAL_SCROLL',
            'horizontalScroll': 'HORIZONTAL_SCROLL',
            'wave': 'WAVE',
            'doubleWave': 'DOUBLE_WAVE',
            'gear6': 'GEAR_6',
            'gear9': 'GEAR_9',
            'funnel': 'FUNNEL',
            'mathPlus': 'MATH_PLUS',
            'mathMinus': 'MATH_MINUS',
            'mathMultiply': 'MATH_MULTIPLY',
            'mathDivide': 'MATH_DIVIDE',
            'mathEqual': 'MATH_EQUAL',
            'mathNotEqual': 'MATH_NOT_EQUAL',
            'homePlate': 'PENTAGON',
            'chevron': 'CHEVRON',
            'pieWedge': 'PIE_WEDGE',
            'irregularSeal1': 'EXPLOSION1',
            'irregularSeal2': 'EXPLOSION2',
            'cloud': 'CLOUD',
            'cloudCallout': 'CLOUD_CALLOUT',
            'callout1': 'RECTANGULAR_CALLOUT',
            'callout2': 'ROUNDED_RECTANGULAR_CALLOUT',
            'callout3': 'OVAL_CALLOUT',
            'wedgeEllipseCallout': 'OVAL_CALLOUT',
            'wedgeRectCallout': 'RECTANGULAR_CALLOUT',
            'wedgeRRectCallout': 'ROUNDED_RECTANGULAR_CALLOUT',
            'borderCallout1': 'LINE_CALLOUT_1',
            'borderCallout2': 'LINE_CALLOUT_2',
            'borderCallout3': 'LINE_CALLOUT_3',
            'accentCallout1': 'LINE_CALLOUT_1_ACCENT_BAR',
            'accentCallout2': 'LINE_CALLOUT_2_ACCENT_BAR',
            'accentCallout3': 'LINE_CALLOUT_3_ACCENT_BAR',
            'callout90': 'LINE_CALLOUT_4',
            'accentCallout90': 'LINE_CALLOUT_4_ACCENT_BAR',
            'borderCallout90': 'LINE_CALLOUT_4_BORDER_AND_ACCENT_BAR',
        }
        
        return geometry_mapping.get(geom_name, geom_name.upper())
    
    def get_shape_type_enum_name(self, shape_type_value: int) -> str:
        """Convert MSO_SHAPE_TYPE enum value to consistent name format"""
        # Map MSO_SHAPE_TYPE enum values to their names
        shape_type_names = {
            -2: 'MIXED',
            1: 'AUTO_SHAPE',
            2: 'CALLOUT',
            3: 'CHART', 
            4: 'COMMENT',
            5: 'FREEFORM',
            6: 'GROUP',
            7: 'EMBEDDED_OLE_OBJECT',
            8: 'FORM_CONTROL',
            9: 'LINE',
            10: 'LINKED_OLE_OBJECT',
            11: 'LINKED_PICTURE',
            12: 'OLE_CONTROL_OBJECT',
            13: 'PICTURE',
            14: 'PLACEHOLDER',
            15: 'TEXT_EFFECT',
            16: 'MEDIA',
            17: 'TEXT_BOX',
            18: 'SCRIPT_ANCHOR',
            19: 'TABLE',
            20: 'CANVAS',
            21: 'DIAGRAM',
            22: 'INK',
            23: 'INK_COMMENT',
            24: 'IGX_GRAPHIC',
            26: 'WEB_VIDEO',
            27: 'CONTENT_APP',
            28: 'GRAPHIC',
            29: 'LINKED_GRAPHIC',
            30: '3D_MODEL',
            31: 'LINKED_3D_MODEL',
        }
        
        return shape_type_names.get(shape_type_value, f'UNKNOWN_{shape_type_value}')

    def extract_chart_data(self, chart) -> Dict[str, Any]:
        """Extract chart data and properties"""
        chart_info = {
            'chart_type': str(chart.chart_type) if hasattr(chart, 'chart_type') else None,
            'has_title': hasattr(chart, 'chart_title') and chart.chart_title.has_text_frame,
            'title': chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and chart.chart_title.has_text_frame else None,
            'categories': [],
            'series': [],
            'has_legend': hasattr(chart, 'has_legend') and chart.has_legend
        }
        
        try:
            # Extract categories
            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    chart_info['categories'] = [cat for cat in plot.categories]
                
                # Extract series data
                if hasattr(plot, 'series'):
                    for series in plot.series:
                        series_info = {
                            'name': series.name if hasattr(series, 'name') else None,
                            'values': [val for val in series.values] if hasattr(series, 'values') else []
                        }
                        chart_info['series'].append(series_info)
                        
        except Exception as e:
            chart_info['error'] = f"Could not extract chart data: {str(e)}"
            
        return chart_info
    
    def extract_table_data(self, table) -> Dict[str, Any]:
        """Extract table data and properties"""
        table_info = {
            'rows': table.rows.__len__() if hasattr(table, 'rows') else 0,
            'columns': table.columns.__len__() if hasattr(table, 'columns') else 0,
            'data': []
        }
        
        try:
            # Extract table cell data
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text if hasattr(cell, 'text') else ''
                    row_data.append(cell_text)
                table_info['data'].append(row_data)
                
        except Exception as e:
            table_info['error'] = f"Could not extract table data: {str(e)}"
            
        return table_info
    
    def extract_placeholder_info(self, shape) -> Dict[str, Any]:
        """Extract placeholder-specific information"""
        placeholder_info = {}
        
        try:
            if hasattr(shape, 'placeholder_format'):
                placeholder_info['placeholder_type'] = str(shape.placeholder_format.type)
                placeholder_info['placeholder_idx'] = shape.placeholder_format.idx if hasattr(shape.placeholder_format, 'idx') else None
                
        except Exception as e:
            placeholder_info['error'] = f"Could not extract placeholder info: {str(e)}"
            
        return placeholder_info
    
    def extract_shapes(self) -> List[Dict[str, Any]]:
        """Extract shape information from all slides"""
        shapes_data = []
        
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_shapes = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_info = {
                    'slide_index': slide_idx,
                    'shape_index': shape_idx,
                    'shape_id': shape.shape_id,
                    'name': shape.name,
                    'shape_type': self.get_auto_shape_type(shape),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                }
                
                # Add text content if available
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape_info['text'] = shape.text_frame.text
                
                # Extract chart data for chart shapes
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    if hasattr(shape, 'chart'):
                        shape_info['chart_data'] = self.extract_chart_data(shape.chart)
                
                # Extract table data for table shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    if hasattr(shape, 'table'):
                        shape_info['table_data'] = self.extract_table_data(shape.table)
                
                # Extract placeholder information
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    shape_info['placeholder_info'] = self.extract_placeholder_info(shape)
                
                # Add detailed fill properties
                if hasattr(shape, 'fill'):
                    shape_info['has_fill'] = True
                    fill_info = self.extract_fill_properties(shape.fill)
                    shape_info['fill'] = fill_info
                
                # Add detailed line properties
                if hasattr(shape, 'line'):
                    shape_info['has_line'] = True
                    line_info = self.extract_line_properties(shape.line)
                    shape_info['line'] = line_info
                
                slide_shapes.append(shape_info)
            
            shapes_data.append({
                'slide_index': slide_idx,
                'shapes': slide_shapes
            })
        
        return shapes_data
    
    def extract_layouts(self) -> List[Dict[str, Any]]:
        """Extract layout information from the presentation"""
        layouts_data = []
        
        for layout_idx, layout in enumerate(self.presentation.slide_layouts):
            layout_info = {
                'layout_index': layout_idx,
                'name': layout.name,
                'placeholders': []
            }
            
            # Extract placeholder information
            for placeholder in layout.placeholders:
                placeholder_info = {
                    'placeholder_format': str(placeholder.placeholder_format.type),
                    'name': placeholder.name,
                    'left': placeholder.left,
                    'top': placeholder.top,
                    'width': placeholder.width,
                    'height': placeholder.height,
                }
                layout_info['placeholders'].append(placeholder_info)
            
            layouts_data.append(layout_info)
        
        return layouts_data
    
    def extract_theme(self) -> Dict[str, Any]:
        """Extract theme information from the presentation"""
        theme_data = {
            'slide_master': {},
            'color_scheme': {},
            'font_scheme': {}
        }
        
        # Extract slide master information
        slide_master = self.presentation.slide_master
        theme_data['slide_master'] = {
            'name': slide_master.name if hasattr(slide_master, 'name') else 'Unknown',
            'background': str(slide_master.background.fill.type) if slide_master.background.fill else None,
        }
        
        # Extract theme colors
        try:
            theme_part = self.presentation.part.theme_part
            if theme_part:
                theme_data['theme_name'] = theme_part.name if hasattr(theme_part, 'name') else 'Default Theme'
        except:
            theme_data['theme_name'] = 'Default Theme'
        
        # Extract color scheme information
        try:
            color_scheme = slide_master.theme.color_scheme
            theme_colors = {}
            for i, color in enumerate(color_scheme):
                theme_colors[f'color_{i}'] = {
                    'rgb': str(color.rgb) if hasattr(color, 'rgb') else None,
                    'type': str(color.color_type) if hasattr(color, 'color_type') else None
                }
            theme_data['color_scheme'] = theme_colors
        except:
            theme_data['color_scheme'] = {'error': 'Could not extract color scheme'}
        
        # Extract font scheme information
        try:
            font_scheme = slide_master.theme.font_scheme
            theme_data['font_scheme'] = {
                'major_font': {
                    'latin': font_scheme.major_font.latin if hasattr(font_scheme.major_font, 'latin') else None,
                    'ea': font_scheme.major_font.ea if hasattr(font_scheme.major_font, 'ea') else None,
                    'cs': font_scheme.major_font.cs if hasattr(font_scheme.major_font, 'cs') else None,
                },
                'minor_font': {
                    'latin': font_scheme.minor_font.latin if hasattr(font_scheme.minor_font, 'latin') else None,
                    'ea': font_scheme.minor_font.ea if hasattr(font_scheme.minor_font, 'ea') else None,
                    'cs': font_scheme.minor_font.cs if hasattr(font_scheme.minor_font, 'cs') else None,
                }
            }
        except:
            theme_data['font_scheme'] = {'error': 'Could not extract font scheme'}
        
        return theme_data
    
    def save_to_json(self, data: Any, output_file: str):
        """Save data to JSON file"""
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        print(f"Data saved to: {output_file}")


def main():
    parser = argparse.ArgumentParser(description='Extract shapes, layouts, and theme from PowerPoint files')
    parser.add_argument('input_file', help='Path to the PowerPoint file (.ppt or .pptx)')
    parser.add_argument('--output-dir', default='.', help='Output directory for JSON files (default: current directory)')
    
    args = parser.parse_args()
    
    # Validate input file
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: File '{args.input_file}' does not exist.")
        sys.exit(1)
    
    if not input_path.suffix.lower() in ['.ppt', '.pptx']:
        print(f"Error: File '{args.input_file}' is not a PowerPoint file (.ppt or .pptx).")
        sys.exit(1)
    
    # Create output directory if it doesn't exist
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Extract base filename for output files
    base_name = input_path.stem
    
    try:
        # Initialize extractor
        extractor = PPTExtractor(args.input_file)
        
        # Extract shapes
        print("Extracting shapes...")
        shapes_data = extractor.extract_shapes()
        shapes_output = output_dir / f"{base_name}_shapes.json"
        extractor.save_to_json(shapes_data, shapes_output)
        
        # Extract layouts
        print("Extracting layouts...")
        layouts_data = extractor.extract_layouts()
        layouts_output = output_dir / f"{base_name}_layouts.json"
        extractor.save_to_json(layouts_data, layouts_output)
        
        # Extract theme
        print("Extracting theme...")
        theme_data = extractor.extract_theme()
        theme_output = output_dir / f"{base_name}_theme.json"
        extractor.save_to_json(theme_data, theme_output)
        
        print(f"\nExtraction completed successfully!")
        print(f"Output files created in: {output_dir}")
        
    except Exception as e:
        print(f"Error processing PowerPoint file: {str(e)}")
        sys.exit(1)


if __name__ == "__main__":
    main()