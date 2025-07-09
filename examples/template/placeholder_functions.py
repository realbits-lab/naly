#!/usr/bin/env python3
"""
Shape creation functions for different PowerPoint placeholder types.
Creates shape instances that simulate different placeholder types.
"""

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_shape(slide, left=Inches(1), top=Inches(0.5), width=Inches(8), height=Inches(1.5)):
    """Create a title text shape."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Title Text"

    # Format as title
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    font = paragraph.font
    font.size = Pt(44)
    font.bold = True

    return textbox


def create_body_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Create a body text shape."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Body content\n• Bullet point 1\n• Bullet point 2"

    # Format as body text
    paragraph = text_frame.paragraphs[0]
    font = paragraph.font
    font.size = Pt(18)

    return textbox


def create_center_title_shape(slide, left=Inches(1), top=Inches(2.5), width=Inches(8), height=Inches(2)):
    """Create a center title shape."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Center Title"

    # Format as center title
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    font = paragraph.font
    font.size = Pt(36)
    font.bold = True

    return textbox


def create_subtitle_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(1)):
    """Create a subtitle shape."""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Subtitle Text"

    # Format as subtitle
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    font = paragraph.font
    font.size = Pt(24)

    return textbox


def create_object_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Create an object shape (rectangle)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)

    # Add text to the shape
    text_frame = shape.text_frame
    text_frame.text = "Object Content"

    return shape


def create_chart_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Create a chart placeholder shape."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # Create sample chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Series 1', (20, 30, 25, 35))

    # Add chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )

    return chart


def create_table_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Create a table shape."""
    rows, cols = 3, 3
    table = slide.shapes.add_table(rows, cols, left, top, width, height)

    # Add sample data
    table.table.cell(0, 0).text = "Header 1"
    table.table.cell(0, 1).text = "Header 2"
    table.table.cell(0, 2).text = "Header 3"

    for row in range(1, rows):
        for col in range(cols):
            table.table.cell(row, col).text = f"Cell {row},{col}"

    return table


def create_clip_art_shape(slide, left=Inches(1), top=Inches(2), width=Inches(4), height=Inches(4)):
    """Create a clip art placeholder shape (using a shape)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)

    # Add text to indicate it's for clip art
    text_frame = shape.text_frame
    text_frame.text = "Clip Art"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    return shape


def create_diagram_shape(slide, left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4)):
    """Create a diagram (SmartArt) placeholder shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    # Add text to indicate it's for diagram
    text_frame = shape.text_frame
    text_frame.text = "Diagram (SmartArt)"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    return shape


def create_media_shape(slide, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(4)):
    """Create a media placeholder shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)

    # Add text to indicate it's for media
    text_frame = shape.text_frame
    text_frame.text = "Media Content"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    # Set fill color to indicate media
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray

    return shape


def create_picture_shape(slide, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(4)):
    """Create a picture placeholder shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)

    # Add text to indicate it's for picture
    text_frame = shape.text_frame
    text_frame.text = "Picture"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER

    # Set border to indicate picture placeholder
    line = shape.line
    line.color.rgb = RGBColor(100, 100, 100)
    line.width = Pt(2)

    return shape


def create_all_placeholder_shapes(slide):
    """Create all placeholder shape types on a slide."""
    shapes = {
        'title': create_title_shape(slide, Inches(1), Inches(0.5), Inches(8), Inches(1)),
        'body': create_body_shape(slide, Inches(1), Inches(2), Inches(4), Inches(4)),
        'center_title': create_center_title_shape(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.5)),
        'subtitle': create_subtitle_shape(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8)),
        'object': create_object_shape(slide, Inches(5.5), Inches(2), Inches(3), Inches(2)),
        'chart': create_chart_shape(slide, Inches(1), Inches(2), Inches(4), Inches(3)),
        'table': create_table_shape(slide, Inches(5.5), Inches(2), Inches(3), Inches(3)),
        'clip_art': create_clip_art_shape(slide, Inches(1), Inches(5.5), Inches(2), Inches(1.5)),
        'diagram': create_diagram_shape(slide, Inches(3.5), Inches(5.5), Inches(2.5), Inches(1.5)),
        'media': create_media_shape(slide, Inches(6.5), Inches(5.5), Inches(2.5), Inches(1.5)),
        'picture': create_picture_shape(slide, Inches(1), Inches(6), Inches(2), Inches(1)),
    }

    return shapes
