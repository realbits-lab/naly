#!/usr/bin/env python3
"""
Generate a PowerPoint file with custom slide using available layouts.
"""

import os
from pptx import Presentation
from pptx.slide import Slide, SlideLayout
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from placeholder_functions import (
    create_title_shape,
    create_body_shape,
    create_subtitle_shape,
    create_chart_shape,
    create_table_shape,
    create_picture_shape,
    create_all_placeholder_shapes
)


def generate_default_powerpoint():
    """Generate a default PowerPoint presentation with python-pptx library."""
    # Check if blank.pptx exists in the current directory
    blank_path = os.path.join(os.path.dirname(__file__), "blank.pptx")
    if os.path.exists(blank_path):
        print(f"Loading template from: {blank_path}")
        prs = Presentation(blank_path)
    else:
        # Create a blank presentation if blank.pptx doesn't exist
        print("Creating new blank presentation")
        prs = Presentation()

    # Show available layouts
    print(f"Available layouts: {len(prs.slide_layouts)}")
    for i, layout in enumerate(prs.slide_layouts):
        print(
            f"  Layout {i}: {layout.name} ({len(layout.placeholders)} placeholders)")

    # Add a slide using an existing layout with placeholders
    add_custom_layout(prs)

    # Add slides with custom shapes
    add_slide_with_title_and_content(prs)
    add_slide_with_title_and_chart(prs)
    add_slide_with_title_and_table(prs)
    add_slide_with_all_shapes(prs)

    return prs


def add_custom_layout(prs):
    """Add a slide using an existing layout and work with its placeholders."""
    # Try to find a layout with placeholders
    layout = None
    for i, slide_layout in enumerate(prs.slide_layouts):
        print(
            f"Layout {i}: {slide_layout.name} - {len(slide_layout.placeholders)} placeholders")
        if len(slide_layout.placeholders) > 0:
            layout = slide_layout
            break

    if layout is None:
        layout = prs.slide_layouts[0]
        print(f"Using default layout: {layout.name}")

    slide = prs.slides.add_slide(layout)

    # Work with existing placeholders
    print(f"Slide has {len(slide.placeholders)} placeholders")

    for placeholder in slide.placeholders:
        idx = placeholder.placeholder_format.idx
        ph_type = placeholder.placeholder_format.type
        print(f"Placeholder idx: {idx}, type: {ph_type}")

        # Add content based on placeholder type
        if ph_type == PP_PLACEHOLDER.TITLE:
            placeholder.text = "Layout Title"
        elif ph_type == PP_PLACEHOLDER.BODY:
            placeholder.text = "Layout Body Content\n• Bullet 1\n• Bullet 2"
        elif ph_type == PP_PLACEHOLDER.OBJECT:
            placeholder.text = "Object Content"

    # If no placeholders exist, create regular shapes as fallback
    if len(slide.placeholders) == 0:
        print("No placeholders found, creating regular shapes")
        title_shape = create_title_shape(slide, Inches(
            1), Inches(0.5), Inches(8), Inches(1.5))
        content_shape = create_body_shape(
            slide, Inches(1), Inches(2), Inches(8), Inches(4))

    print(f"Added slide using layout: {layout.name}")
    return slide


def add_slide_with_custom_shapes(prs, shape_types=['title', 'body']):
    """Add a slide with custom shapes simulating different placeholder types."""
    # Use blank layout or first available layout
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    created_shapes = []

    for shape_type in shape_types:
        if shape_type == 'title':
            shape = create_title_shape(slide)
            created_shapes.append(('title', shape))
        elif shape_type == 'body':
            shape = create_body_shape(slide)
            created_shapes.append(('body', shape))
        elif shape_type == 'subtitle':
            shape = create_subtitle_shape(slide)
            created_shapes.append(('subtitle', shape))
        elif shape_type == 'chart':
            shape = create_chart_shape(slide)
            created_shapes.append(('chart', shape))
        elif shape_type == 'table':
            shape = create_table_shape(slide)
            created_shapes.append(('table', shape))
        elif shape_type == 'picture':
            shape = create_picture_shape(slide)
            created_shapes.append(('picture', shape))

    print(
        f"Created slide with {len(created_shapes)} custom shapes: {[s[0] for s in created_shapes]}")
    return slide, created_shapes


def add_slide_with_title_and_content(prs):
    """Add a slide with title and body content shapes."""
    return add_slide_with_custom_shapes(prs, ['title', 'body'])


def add_slide_with_title_and_chart(prs):
    """Add a slide with title and chart shapes."""
    return add_slide_with_custom_shapes(prs, ['title', 'chart'])


def add_slide_with_title_and_table(prs):
    """Add a slide with title and table shapes."""
    return add_slide_with_custom_shapes(prs, ['title', 'table'])


def add_slide_with_all_shapes(prs):
    """Add a slide with all available shape types."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    shapes = create_all_placeholder_shapes(slide)

    print(
        f"Created slide with all {len(shapes)} shape types: {list(shapes.keys())}")
    return slide, shapes


def save_powerpoint(prs, filename="default_template.pptx"):
    """Save the PowerPoint presentation to a file."""
    filepath = os.path.join(os.path.dirname(__file__), filename)

    # Check if file exists
    if os.path.exists(filepath):
        print(f"File {filename} already exists. Replacing it.")
        os.remove(filepath)

    # Save the presentation
    prs.save(filepath)
    print(f"PowerPoint saved as: {filepath}")

    # Show file info
    file_size = os.path.getsize(filepath)
    print(f"File size: {file_size} bytes")


def main():
    """Main function to generate PowerPoint with custom slide."""
    print("Generating PowerPoint presentation...")

    # Generate the presentation
    prs = generate_default_powerpoint()

    # Save to file
    save_powerpoint(prs)

    print("\n=== Summary ===")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Total slide layouts: {len(prs.slide_layouts)}")

    # Show information about the layouts
    if len(prs.slide_layouts) > 0:
        layout = prs.slide_layouts[0]
        print(f"First layout name: {layout.name}")
        print(f"First layout placeholders: {len(layout.placeholders)}")

    print("Presentation created successfully!")


if __name__ == "__main__":
    main()
